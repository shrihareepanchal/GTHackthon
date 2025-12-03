from pptx import Presentation
from pptx.util import Inches, Pt

def export_ppt(title, summary, metrics, images, anomalies, recommendations, outpath):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = ""

    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "AI Summary"
    tb = summary_slide.placeholders[1].text_frame
    for line in summary.split("\n"):
        p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    metrics_slide = prs.slides.add_slide(prs.slide_layouts[1])
    metrics_slide.shapes.title.text = "Key Metrics"
    tb = metrics_slide.placeholders[1].text_frame
    for k, v in metrics.items():
        p = tb.add_paragraph()
        p.text = f"{k}: {v}"
        p.font.size = Pt(14)

    anomaly_slide = prs.slides.add_slide(prs.slide_layouts[1])
    anomaly_slide.shapes.title.text = "Detected Anomalies"
    tb = anomaly_slide.placeholders[1].text_frame
    if len(anomalies) == 0:
        p = tb.add_paragraph()
        p.text = "No anomalies detected."
        p.font.size = Pt(14)
    else:
        for a in anomalies:
            p = tb.add_paragraph()
            p.text = a
            p.font.size = Pt(14)

    for img in images:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Visualization"
        s.shapes.add_picture(img, Inches(1), Inches(1), width=Inches(6))

    rec_slide = prs.slides.add_slide(prs.slide_layouts[1])
    rec_slide.shapes.title.text = "Recommendations"
    tb = rec_slide.placeholders[1].text_frame

    rec_text = "\n".join(recommendations)
    for line in rec_text.split("\n"):
        p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    prs.save(outpath)
    return outpath
